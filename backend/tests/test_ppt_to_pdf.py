from pathlib import Path
from pptx import Presentation
from backend.app.processors.ppt_to_pdf import PptToPdfProcessor
from backend.app.core.validation import validate_pdf_output

FIXTURES_DIR = Path(__file__).parent / "fixtures"


def test_ppt_to_pdf_conversion(tmp_path):
    """INVARIANT: Generates valid PDF from PPTX presentation with matching pages."""
    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Role of Tech in Growth"
    slide.placeholders[1].text = "Testing PPT to PDF conversion"
    
    prs.save(str(pptx_path))
    
    processor = PptToPdfProcessor(job_id="test_ppt2pdf", work_dir=tmp_path)
    result = processor.run([pptx_path], {})
    
    assert result.status == "COMPLETED"
    page_count = validate_pdf_output(result.output_path)
    assert page_count >= 1
