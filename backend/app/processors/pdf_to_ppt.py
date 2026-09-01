import os
from pathlib import Path
from typing import Any, Dict

import pymupdf as fitz
from pptx import Presentation
from pptx.util import Inches, Pt

from backend.app.core.errors import PDFBoltError, OutputValidationError
from backend.app.core.validation import validate_pptx_output
from backend.app.processors.base import BaseProcessor
from backend.app.core.logging import logger


class PdfToPptProcessor(BaseProcessor):
    """
    Direct PDF -> PPTX Conversion Engine using PyMuPDF and python-pptx.
    Renders each PDF page at high fidelity (200 DPI) and embeds it matching exact page dimensions.
    """

    operation = "pdf-to-ppt"
    input_formats = [".pdf"]
    output_format = ".pptx"

    def process(self, input_files: Any, options: Any = None) -> Any:
        if isinstance(input_files, (bytes, bytearray)):
            return self.process_bytes(input_files, str(options or "doc.pdf"))

        if not input_files:
            raise PDFBoltError("NO_FILES_PROVIDED", "No input PDF provided for conversion.")

        input_pdf = Path(input_files[0])
        if not input_pdf.exists():
            raise PDFBoltError("FILE_NOT_FOUND", f"Input PDF file not found: {input_pdf}")

        output_path = self.output_dir / f"{self.job_id}.pptx"

        prs = Presentation()

        try:
            logger.info(f"Converting '{input_pdf}' to PPTX with PyMuPDF and python-pptx...")
            doc = fitz.open(str(input_pdf))

            if doc.page_count == 0:
                doc.close()
                raise PDFBoltError("INVALID_DOCUMENT", "The PDF document is empty.")

            # Get dimensions of the first page to set slide size
            first_page = doc[0]
            rect = first_page.rect
            slide_width_pt = rect.width
            slide_height_pt = rect.height

            # Set presentation slide size to match PDF page size
            prs.slide_width = Pt(slide_width_pt)
            prs.slide_height = Pt(slide_height_pt)

            for page_num in range(doc.page_count):
                page = doc.load_page(page_num)  # Load page

                # Render page to PNG image at 200 DPI (200 / 72 zoom matrix)
                pix = page.get_pixmap(matrix=fitz.Matrix(200 / 72, 200 / 72))
                img_path = self.temp_dir / f"page_{page_num + 1}.png"
                pix.save(str(img_path))

                # Add a blank slide to the presentation
                blank_slide_layout = prs.slide_layouts[6]  # Blank slide layout
                slide = prs.slides.add_slide(blank_slide_layout)

                # Add image to slide, fitting it to the slide dimensions
                slide.shapes.add_picture(
                    str(img_path),
                    Inches(0),
                    Inches(0),
                    width=prs.slide_width,
                    height=prs.slide_height
                )

                # Clean up the temporary image file
                if img_path.exists():
                    os.remove(str(img_path))

            doc.close()
            prs.save(str(output_path))
            logger.info(f"Successfully converted '{input_pdf}' to '{output_path}'")
        except Exception as e:
            if isinstance(e, PDFBoltError):
                raise
            logger.error(f"Error during PDF to PPT conversion: {e}")
            raise PDFBoltError("CONVERSION_FAILED", f"PDF to PPT conversion failed: {e}")

        if not output_path.exists() or output_path.stat().st_size == 0:
            raise OutputValidationError("Failed to generate a valid PowerPoint (.pptx) presentation.")

        validate_pptx_output(output_path)

        self.metrics = {
            "format": "pptx",
            "quality_status": "passed",
            "quality_score": 100,
            "status": "success"
        }

        return output_path

    def process_bytes(self, content: bytes, filename: str) -> tuple[bytes, str, Dict[str, Any]]:
        temp_in = self.temp_dir / "in.pdf"
        with open(temp_in, "wb") as f:
            f.write(content)

        out_path = self.process([temp_in], self.settings)
        with open(out_path, "rb") as f:
            out_bytes = f.read()

        metrics = dict(getattr(self, "metrics", {}))
        metrics["original_size_bytes"] = len(content)
        metrics["output_size_bytes"] = len(out_bytes)
        metrics["format"] = "pptx"
        metrics["quality_status"] = "passed"
        metrics["quality_score"] = 100

        return out_bytes, "presentation.pptx", metrics


PDFToPPTProcessor = PdfToPptProcessor
PdfToPptProcessor = PdfToPptProcessor
